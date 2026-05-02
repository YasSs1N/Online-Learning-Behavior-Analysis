"""build_slides.py — assembles the 12-slide PowerPoint deck."""
from pathlib import Path
import pandas as pd
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parents[1]
RESULTS = ROOT / "results"
FIGS = ROOT / "figures"
OUT = ROOT / "slides" / "Online_Learning_Behavior_Analysis.pptx"
OUT.parent.mkdir(parents=True, exist_ok=True)

PRIMARY = RGBColor(0x1F, 0x3A, 0x68)
ACCENT = RGBColor(0x4C, 0x8F, 0xD9)


def title(s, text, size=30):
    tb = s.shapes.add_textbox(Cm(1.0), Cm(0.6), Cm(23.5), Cm(2.0))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    for r in p.runs:
        r.font.size = Pt(size); r.font.bold = True
        r.font.color.rgb = PRIMARY; r.font.name = "Calibri"
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(1.0), Cm(2.5), Cm(4.0), Cm(0.15))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()


def bullets(s, items, top=Cm(3.2), height=Cm(13), size=17):
    box = s.shapes.add_textbox(Cm(1.0), top, Cm(23.5), height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.space_after = Pt(8)
        for r in p.runs:
            r.font.size = Pt(size); r.font.name = "Calibri"
            r.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


def img(s, name, **kw):
    f = FIGS / name
    if f.exists():
        s.shapes.add_picture(str(f), **kw)


def footer(s, idx, total):
    tb = s.shapes.add_textbox(Cm(1.0), Cm(18.5), Cm(23.5), Cm(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = f"Online Learning Behavior Analysis  -  {idx}/{total}"
    p.alignment = PP_ALIGN.RIGHT
    for r in p.runs:
        r.font.size = Pt(9); r.font.color.rgb = RGBColor(0x7F, 0x7F, 0x7F)


prs = Presentation()
prs.slide_width = Cm(25.4); prs.slide_height = Cm(19.05)
blank = prs.slide_layouts[6]

# Slide 1 — Title
s = prs.slides.add_slide(blank)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = PRIMARY; bg.line.fill.background()
tb = s.shapes.add_textbox(Cm(1.0), Cm(6.0), Cm(23.5), Cm(4.0))
p = tb.text_frame.paragraphs[0]; p.text = "Online Learning Behavior Analysis"
for r in p.runs:
    r.font.size = Pt(40); r.font.bold = True
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF); r.font.name = "Calibri"
sub = s.shapes.add_textbox(Cm(1.0), Cm(9.5), Cm(23.5), Cm(2.0))
p = sub.text_frame.paragraphs[0]; p.text = "Project #6 — Data Mining Course"
for r in p.runs:
    r.font.size = Pt(20); r.font.italic = True
    r.font.color.rgb = RGBColor(0xCF, 0xDB, 0xEC)
au = s.shapes.add_textbox(Cm(1.0), Cm(15.5), Cm(23.5), Cm(2.0))
p = au.text_frame.paragraphs[0]; p.text = "Yassin Mekawy  |  April 2026"
for r in p.runs:
    r.font.size = Pt(16); r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# Slide 2 — Why two datasets
s = prs.slides.add_slide(blank); title(s, "1. Datasets — one per task")
bullets(s, [
    "OULAD (Open University Learning Analytics) — 32k+ student × module-presentation records over 7 modules. Used for Task 1: real co-enrolment data for FP-Growth & PageRank.",
    "Coursera Courses Dataset 2021 — 3,522 courses with name, description (~1,160 chars), and skills. Used for Task 2: real course descriptions for BERT.",
    "No single public dataset has both rich enrolments and full descriptions — the two-dataset approach is the methodologically honest answer.",
])

# Slide 3 — Module mapping
s = prs.slides.add_slide(blank); title(s, "2. OULAD module mapping")
bullets(s, [
    "OULAD modules are anonymised AAA–GGG. Kuzilek et al. (2017) disclose 4 STEM + 3 Social Sciences.",
    "Subject-area names inferred from enrolment, gender, and assessment signature (transparent in the report).",
    "Single source of truth: data/processed/oulad_module_mapping.xlsx — every script reads this file.",
    "Visualisations always use the readable name, never the code.",
])
img(s, "01_oulad_overview.png", left=Cm(1.0), top=Cm(10.0), width=Cm(23.5))

# Slide 4 — Task 1 association rules
s = prs.slides.add_slide(blank); title(s, "3. Task 1 — FP-Growth + Apriori")
rules = pd.read_csv(RESULTS / "association_rules.csv")
bullets(s, [
    "Each student is a transaction = the modules they enrolled in (real multi-enrolments, no synthesis).",
    f"FP-Growth at min_support=0.005 → {len(rules)} rules with lift ≥ 1.",
    "Strongest co-enrolment: Mathematics ↔ Computing & IT, lift ≈ 2.52.",
    "This is the \"common course combinations\" the rubric asks for — discovered, not invented.",
], top=Cm(3.2), height=Cm(7))
img(s, "02_association_rules.png", left=Cm(1.0), top=Cm(11.0), width=Cm(23.5))

# Slide 5 — Task 1 PageRank
s = prs.slides.add_slide(blank); title(s, "4. Task 1 — PageRank on the course graph")
bullets(s, [
    "Nodes = modules. Edges weighted by # students enrolled in both.",
    "PageRank α=0.85 on the weighted graph.",
    "Computing & IT is the central module — gateway between four STEM modules.",
    "Smaller-but-connected beats larger-but-isolated.",
], top=Cm(3.2), height=Cm(5))
img(s, "03_course_graph_pagerank.png", left=Cm(1.0), top=Cm(8.5), width=Cm(13.0))
img(s, "04_pagerank_bar.png", left=Cm(14.5), top=Cm(8.5), width=Cm(10.5))

# Slide 6 — HITS
s = prs.slides.add_slide(blank); title(s, "5. Task 1 — HITS hubs & authorities")
bullets(s, [
    "HITS gives each module two scores: Hub (sends edges to authoritative neighbours) and Authority (receives them).",
    "Computing & IT is both top hub AND top authority — confirming PageRank.",
    "Social Sciences modules score near zero on both axes (peripheral).",
], top=Cm(3.2), height=Cm(5))
img(s, "05_hits.png", left=Cm(2.0), top=Cm(9.0), width=Cm(21.5))

# Slide 7 — Coursera + BERT setup
s = prs.slides.add_slide(blank); title(s, "6. Task 2 — Coursera + BERT setup")
bullets(s, [
    "3,416 unique courses, descriptions averaging 1,160 characters.",
    "Topic-tagged into Computing & IT, Mathematics, Engineering, Humanities, Education, Social Sciences, Business, Health & Medicine, Other.",
    "Model A: sentence-transformers all-MiniLM-L6-v2 (BERT-family, fine-tuned for sentence similarity).",
    "Model B: bert-base-uncased + manual mean-pooling (raw BERT, literal interpretation).",
    "Cosine similarity over normalised embeddings; top-10 neighbours per course.",
])

# Slide 8 — BERT topic heatmap
s = prs.slides.add_slide(blank); title(s, "7. BERT — topic-level similarity")
bullets(s, [
    "Aggregating per-course similarities by topic shows strong diagonal structure.",
    "Within-topic similarity dominates between-topic — BERT correctly groups by content.",
], top=Cm(3.2), height=Cm(3.5))
img(s, "06_bert_topic_heatmap.png", left=Cm(3.0), top=Cm(7.0), width=Cm(19.5))

# Slide 9 — BERT comparison
s = prs.slides.add_slide(blank); title(s, "8. BERT — model comparison")
img(s, "07_bert_model_comparison.png", left=Cm(1.0), top=Cm(3.5), width=Cm(23.5))

# Slide 10 — BERT examples
s = prs.slides.add_slide(blank); title(s, "9. BERT — sample neighbours")
nb = pd.read_csv(RESULTS / "bert_minilm_top_neighbours.csv")
top = (nb[nb["rank"] == 1].sort_values("similarity", ascending=False).head(7))
bullets(s, [f"{r.course_name}  →  {r.neighbour_name}  (sim={r.similarity:.2f})"
            for r in top.itertuples()],
        top=Cm(3.2), height=Cm(11), size=14)

# Slide 11 — Recommended paths
s = prs.slides.add_slide(blank); title(s, "10. Recommended learning paths")
bullets(s, [
    "OULAD pattern signal (FP-Growth) — what students DO take together.",
    "PageRank centrality — gateway courses to prioritise.",
    "BERT similarity — what students COULD take next based on content.",
    "Combined into recommended_paths.csv with evidence + rationale per row.",
])

# Slide 12 — Conclusions
s = prs.slides.add_slide(blank); title(s, "11. Conclusions")
bullets(s, [
    "Mathematics ↔ Computing & IT is the strongest real curricular bridge in OULAD (lift > 2.5).",
    "Computing & IT is the platform's most central module by PageRank and HITS.",
    "BERT clusters Coursera courses by content and surfaces course sequels automatically.",
    "Two BERT variants agree at high Pearson correlation — choice of encoder is robust.",
    "Final deliverable: recommended_paths.csv combining pattern + graph + semantic signals.",
])

total = len(prs.slides)
for i, sl in enumerate(prs.slides, start=1):
    if i == 1: continue
    footer(sl, i, total)

prs.save(OUT)
print(f"Wrote {OUT}  ({total} slides)")
